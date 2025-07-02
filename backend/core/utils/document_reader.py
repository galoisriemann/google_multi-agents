"""Document Reader Utility for Multiple File Formats.

This module provides functionality to read PDF, DOCX, TXT, Markdown, CSV, and PPTX files
and convert their content to markdown format for use in agent workflows.
"""

import csv
import logging
from pathlib import Path
from typing import Optional, Dict, Any, List

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

logger = logging.getLogger(__name__)


class DocumentReaderError(Exception):
    """Custom exception for document reading errors."""
    pass


class DocumentReader:
    """Reader class for multiple document formats with markdown output.
    
    Supported formats:
    - PDF (.pdf)
    - Microsoft Word (.docx)
    - Plain Text (.txt)
    - Markdown (.md)
    - CSV (.csv)
    - PowerPoint (.pptx)
    """
    
    def __init__(self, input_directory: Optional[Path] = None):
        """Initialize the document reader.
        
        Args:
            input_directory: Directory containing input documents.
                           Defaults to 'backend/input' if not provided.
        """
        if input_directory is None:
            # Default to backend/input directory
            self.input_directory = Path(__file__).parent.parent.parent / "input"
        else:
            self.input_directory = Path(input_directory)
        
        # Create input directory if it doesn't exist
        self.input_directory.mkdir(parents=True, exist_ok=True)
        
        # Create input_markdown directory for saving markdown versions
        self.input_markdown_directory = self.input_directory.parent / "input_markdown"
        self.input_markdown_directory.mkdir(parents=True, exist_ok=True)
        
        logger.info(f"DocumentReader initialized with input directory: {self.input_directory}")
        logger.info(f"Markdown output directory: {self.input_markdown_directory}")
    
    def read_document(self, filename: str) -> str:
        """Read a document and return its content as markdown.
        
        Args:
            filename: Name of the file to read (with extension)
            
        Returns:
            Document content as markdown string. Empty string if file not found.
            
        Raises:
            DocumentReaderError: If there's an error reading the document
        """
        if not filename or not filename.strip():
            return ""
        
        file_path = self.input_directory / filename
        
        if not file_path.exists():
            logger.warning(f"Document not found: {file_path}")
            return ""
        
        # Check if markdown version already exists
        # Use original filename with .md extension to avoid conflicts between different formats
        markdown_filename = file_path.name + ".md"
        markdown_path = self.input_markdown_directory / markdown_filename
        
        if markdown_path.exists():
            # Return existing markdown version
            logger.info(f"Using cached markdown version: {markdown_path}")
            try:
                with open(markdown_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except Exception as e:
                logger.warning(f"Failed to read cached markdown, regenerating: {e}")
        
        try:
            # Determine file type by extension and read content
            extension = file_path.suffix.lower()
            
            if extension == '.pdf':
                content = self._read_pdf(file_path)
            elif extension == '.docx':
                content = self._read_docx(file_path)
            elif extension == '.txt':
                content = self._read_txt(file_path)
            elif extension == '.md':
                content = self._read_markdown(file_path)
            elif extension == '.csv':
                content = self._read_csv(file_path)
            elif extension == '.pptx':
                content = self._read_pptx(file_path)
            else:
                logger.warning(f"Unsupported file type: {extension}")
                logger.info("Supported formats: .pdf, .docx, .txt, .md, .csv, .pptx")
                return ""
            
            # Save markdown version to input_markdown directory
            if content:
                self._save_markdown_version(content, markdown_path, filename)
            
            return content
                
        except Exception as e:
            error_msg = f"Error reading document {filename}: {str(e)}"
            logger.error(error_msg)
            raise DocumentReaderError(error_msg) from e
    
    def _save_markdown_version(self, content: str, markdown_path: Path, original_filename: str) -> None:
        """Save markdown version of the document to input_markdown directory.
        
        Args:
            content: Markdown content to save
            markdown_path: Path where to save the markdown file
            original_filename: Original filename for logging
        """
        try:
            with open(markdown_path, 'w', encoding='utf-8') as f:
                f.write(content)
            logger.info(f"Saved markdown version: {markdown_path} (from {original_filename})")
        except Exception as e:
            logger.warning(f"Failed to save markdown version for {original_filename}: {e}")
    
    def _read_pdf(self, file_path: Path) -> str:
        """Read PDF file and convert to markdown.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            PDF content as markdown string
            
        Raises:
            DocumentReaderError: If PyPDF2 is not installed or reading fails
        """
        if PyPDF2 is None:
            raise DocumentReaderError(
                "PyPDF2 is not installed. Please install it with: pip install PyPDF2"
            )
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                content_parts = []
                
                # Add metadata header
                content_parts.append(f"# Document: {file_path.name}\n")
                content_parts.append(f"**Source**: PDF file with {len(pdf_reader.pages)} pages\n")
                content_parts.append("---\n")
                
                # Extract text from each page
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            content_parts.append(f"## Page {page_num}\n")
                            content_parts.append(page_text.strip())
                            content_parts.append("\n")
                    except Exception as e:
                        logger.warning(f"Error reading page {page_num}: {e}")
                        content_parts.append(f"## Page {page_num}\n")
                        content_parts.append("*[Error reading this page]*\n")
                
                result = "\n".join(content_parts)
                logger.info(f"Successfully read PDF: {file_path.name} ({len(pdf_reader.pages)} pages)")
                return result
                
        except Exception as e:
            raise DocumentReaderError(f"Failed to read PDF {file_path.name}: {str(e)}") from e
    
    def _read_docx(self, file_path: Path) -> str:
        """Read DOCX file and convert to markdown.
        
        Args:
            file_path: Path to the DOCX file
            
        Returns:
            DOCX content as markdown string
            
        Raises:
            DocumentReaderError: If python-docx is not installed or reading fails
        """
        if Document is None:
            raise DocumentReaderError(
                "python-docx is not installed. Please install it with: pip install python-docx"
            )
        
        try:
            doc = Document(file_path)
            content_parts = []
            
            # Add metadata header
            content_parts.append(f"# Document: {file_path.name}\n")
            content_parts.append(f"**Source**: DOCX file with {len(doc.paragraphs)} paragraphs\n")
            content_parts.append("---\n")
            
            # Extract text from paragraphs
            for i, paragraph in enumerate(doc.paragraphs):
                text = paragraph.text.strip()
                if text:
                    # Simple heuristic for headers (all caps or starts with numbers)
                    if text.isupper() or (text and text[0].isdigit() and '.' in text[:10]):
                        content_parts.append(f"## {text}\n")
                    else:
                        content_parts.append(f"{text}\n")
            
            # Extract text from tables if any
            if doc.tables:
                content_parts.append("\n## Tables\n")
                for table_num, table in enumerate(doc.tables, 1):
                    content_parts.append(f"### Table {table_num}\n")
                    content_parts.append("| " + " | ".join("Column " + str(i+1) for i in range(len(table.rows[0].cells))) + " |\n")
                    content_parts.append("| " + " | ".join("---" for _ in range(len(table.rows[0].cells))) + " |\n")
                    
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip().replace('\n', ' ')
                            row_data.append(cell_text if cell_text else " ")
                        content_parts.append("| " + " | ".join(row_data) + " |\n")
                    content_parts.append("\n")
            
            result = "\n".join(content_parts)
            logger.info(f"Successfully read DOCX: {file_path.name} ({len(doc.paragraphs)} paragraphs)")
            return result
            
        except Exception as e:
            raise DocumentReaderError(f"Failed to read DOCX {file_path.name}: {str(e)}") from e
    
    def _read_txt(self, file_path: Path) -> str:
        """Read TXT file and convert to markdown.
        
        Args:
            file_path: Path to the TXT file
            
        Returns:
            TXT content as markdown string
            
        Raises:
            DocumentReaderError: If reading fails
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            content_parts = []
            
            # Add metadata header
            content_parts.append(f"# Document: {file_path.name}\n")
            content_parts.append(f"**Source**: Plain text file ({len(content)} characters)\n")
            content_parts.append("---\n")
            
            # Add content with basic formatting
            lines = content.split('\n')
            content_parts.append("## Content\n")
            
            for line in lines:
                line = line.strip()
                if line:
                    # Simple heuristic for headers (lines in all caps or with specific patterns)
                    if line.isupper() and len(line) < 100:
                        content_parts.append(f"### {line}\n")
                    elif line.startswith('---') or line.startswith('==='):
                        content_parts.append("---\n")
                    else:
                        content_parts.append(f"{line}\n")
                else:
                    content_parts.append("\n")
            
            result = "\n".join(content_parts)
            logger.info(f"Successfully read TXT: {file_path.name} ({len(content)} characters)")
            return result
            
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    content = file.read()
                logger.warning(f"Read TXT file with latin-1 encoding: {file_path.name}")
                # Format with latin-1 content
                content_parts = [
                    f"# Document: {file_path.name}\n",
                    f"**Source**: Plain text file ({len(content)} characters, latin-1 encoding)\n",
                    "---\n",
                    "## Content\n",
                    content
                ]
                return "\n".join(content_parts)
            except Exception as e:
                raise DocumentReaderError(f"Failed to read TXT {file_path.name} with UTF-8 or latin-1: {str(e)}") from e
        except Exception as e:
            raise DocumentReaderError(f"Failed to read TXT {file_path.name}: {str(e)}") from e
    
    def _read_markdown(self, file_path: Path) -> str:
        """Read Markdown file and return content.
        
        Args:
            file_path: Path to the Markdown file
            
        Returns:
            Markdown content with metadata header
            
        Raises:
            DocumentReaderError: If reading fails
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            content_parts = []
            
            # Add metadata header
            content_parts.append(f"# Document: {file_path.name}\n")
            content_parts.append(f"**Source**: Markdown file ({len(content)} characters)\n")
            content_parts.append("---\n")
            
            # Add original markdown content
            content_parts.append("## Original Content\n")
            content_parts.append(content)
            
            result = "\n".join(content_parts)
            logger.info(f"Successfully read Markdown: {file_path.name} ({len(content)} characters)")
            return result
            
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    content = file.read()
                logger.warning(f"Read Markdown file with latin-1 encoding: {file_path.name}")
                # Format with latin-1 content
                content_parts = [
                    f"# Document: {file_path.name}\n",
                    f"**Source**: Markdown file ({len(content)} characters, latin-1 encoding)\n",
                    "---\n",
                    "## Original Content\n",
                    content
                ]
                return "\n".join(content_parts)
            except Exception as e:
                raise DocumentReaderError(f"Failed to read Markdown {file_path.name} with UTF-8 or latin-1: {str(e)}") from e
        except Exception as e:
            raise DocumentReaderError(f"Failed to read Markdown {file_path.name}: {str(e)}") from e
    
    def _read_csv(self, file_path: Path) -> str:
        """Read CSV file and convert to markdown table.
        
        Args:
            file_path: Path to the CSV file
            
        Returns:
            CSV content as markdown table
            
        Raises:
            DocumentReaderError: If reading fails
        """
        try:
            content_parts = []
            
            # Add metadata header
            content_parts.append(f"# Document: {file_path.name}\n")
            
            with open(file_path, 'r', encoding='utf-8', newline='') as file:
                # Detect delimiter
                sample = file.read(1024)
                file.seek(0)
                sniffer = csv.Sniffer()
                delimiter = sniffer.sniff(sample).delimiter
                
                reader = csv.reader(file, delimiter=delimiter)
                rows = list(reader)
            
            if not rows:
                content_parts.append("**Source**: Empty CSV file\n")
                content_parts.append("---\n")
                content_parts.append("*No data found in CSV file*\n")
                result = "\n".join(content_parts)
                logger.info(f"Successfully read empty CSV: {file_path.name}")
                return result
            
            content_parts.append(f"**Source**: CSV file with {len(rows)} rows and {len(rows[0]) if rows else 0} columns\n")
            content_parts.append("---\n")
            
            # Create markdown table
            content_parts.append("## Data Table\n")
            
            # Header row
            if rows:
                header = rows[0]
                content_parts.append("| " + " | ".join(str(cell).strip() for cell in header) + " |")
                content_parts.append("| " + " | ".join("---" for _ in header) + " |")
                
                # Data rows
                for row in rows[1:]:
                    # Pad row to match header length
                    padded_row = row + [""] * (len(header) - len(row))
                    content_parts.append("| " + " | ".join(str(cell).strip() for cell in padded_row[:len(header)]) + " |")
            
            # Add summary statistics
            content_parts.append(f"\n## Summary\n")
            content_parts.append(f"- **Total Rows**: {len(rows)}")
            content_parts.append(f"- **Total Columns**: {len(rows[0]) if rows else 0}")
            content_parts.append(f"- **Data Rows**: {len(rows) - 1 if len(rows) > 1 else 0}")
            
            result = "\n".join(content_parts)
            logger.info(f"Successfully read CSV: {file_path.name} ({len(rows)} rows)")
            return result
            
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1', newline='') as file:
                    reader = csv.reader(file)
                    rows = list(reader)
                logger.warning(f"Read CSV file with latin-1 encoding: {file_path.name}")
                # Format CSV content with latin-1 encoding
                content_parts = [
                    f"# Document: {file_path.name}\n",
                    f"**Source**: CSV file with {len(rows)} rows (latin-1 encoding)\n",
                    "---\n",
                    "## Data Table\n"
                ]
                
                if rows:
                    header = rows[0]
                    content_parts.append("| " + " | ".join(str(cell).strip() for cell in header) + " |")
                    content_parts.append("| " + " | ".join("---" for _ in header) + " |")
                    
                    for row in rows[1:]:
                        padded_row = row + [""] * (len(header) - len(row))
                        content_parts.append("| " + " | ".join(str(cell).strip() for cell in padded_row[:len(header)]) + " |")
                
                return "\n".join(content_parts)
            except Exception as e:
                raise DocumentReaderError(f"Failed to read CSV {file_path.name} with UTF-8 or latin-1: {str(e)}") from e
        except Exception as e:
            raise DocumentReaderError(f"Failed to read CSV {file_path.name}: {str(e)}") from e
    
    def _read_pptx(self, file_path: Path) -> str:
        """Read PPTX file and convert to markdown.
        
        Args:
            file_path: Path to the PPTX file
            
        Returns:
            PPTX content as markdown string
            
        Raises:
            DocumentReaderError: If python-pptx is not installed or reading fails
        """
        if Presentation is None:
            raise DocumentReaderError(
                "python-pptx is not installed. Please install it with: pip install python-pptx"
            )
        
        try:
            prs = Presentation(file_path)
            content_parts = []
            
            # Add metadata header
            content_parts.append(f"# Document: {file_path.name}\n")
            content_parts.append(f"**Source**: PowerPoint file with {len(prs.slides)} slides\n")
            content_parts.append("---\n")
            
            # Extract content from each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"## Slide {slide_num}\n")
                
                # Extract text from all shapes in the slide
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        # Simple heuristic for titles (usually shorter and on top)
                        if len(text) < 100 and slide_text == []:
                            slide_text.append(f"### {text}\n")
                        else:
                            slide_text.append(f"{text}\n")
                
                if slide_text:
                    content_parts.extend(slide_text)
                else:
                    content_parts.append("*[No text content found on this slide]*\n")
                
                content_parts.append("\n")
            
            # Add summary
            content_parts.append("## Presentation Summary\n")
            content_parts.append(f"- **Total Slides**: {len(prs.slides)}")
            
            result = "\n".join(content_parts)
            logger.info(f"Successfully read PPTX: {file_path.name} ({len(prs.slides)} slides)")
            return result
            
        except Exception as e:
            raise DocumentReaderError(f"Failed to read PPTX {file_path.name}: {str(e)}") from e
    
    def list_available_documents(self) -> Dict[str, Dict[str, Any]]:
        """List all available documents in the input directory.
        
        Supports: PDF, DOCX, TXT, Markdown, CSV, and PPTX files.
        
        Returns:
            Dictionary with filename as key and file info as value
        """
        documents = {}
        supported_extensions = ['.pdf', '.docx', '.txt', '.md', '.csv', '.pptx']
        
        for file_path in self.input_directory.iterdir():
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                documents[file_path.name] = {
                    'path': str(file_path),
                    'type': file_path.suffix.lower(),
                    'size': file_path.stat().st_size,
                    'modified': file_path.stat().st_mtime
                }
        
        return documents
    
    def validate_document(self, filename: str) -> bool:
        """Validate if a document exists and is readable.
        
        Args:
            filename: Name of the file to validate
            
        Returns:
            True if document exists and is supported, False otherwise
        """
        if not filename or not filename.strip():
            return False
        
        file_path = self.input_directory / filename
        
        if not file_path.exists():
            return False
        
        extension = file_path.suffix.lower()
        supported_extensions = ['.pdf', '.docx', '.txt', '.md', '.csv', '.pptx']
        return extension in supported_extensions


# Convenience functions for direct usage
def read_document(filename: str, input_directory: Optional[Path] = None) -> str:
    """Convenience function to read a document.
    
    Args:
        filename: Name of the file to read
        input_directory: Directory containing the file
        
    Returns:
        Document content as markdown string
    """
    reader = DocumentReader(input_directory)
    return reader.read_document(filename)


def list_documents(input_directory: Optional[Path] = None) -> Dict[str, Dict[str, Any]]:
    """Convenience function to list available documents.
    
    Args:
        input_directory: Directory to search for documents
        
    Returns:
        Dictionary of available documents
    """
    reader = DocumentReader(input_directory)
    return reader.list_available_documents() 